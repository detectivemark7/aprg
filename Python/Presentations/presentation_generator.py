# Quick fast guide: https://python-pptx.readthedocs.io/en/latest/user/quickstart.html

from pptx import Presentation


def change_title_slide(presentation):
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides[0]

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = 'This is a title'
    subtitle.text = 'This is a subtitle'


def change_normal_slide(slide):

    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Adding a Bullet Slide'
    text_frame = body_shape.text_frame
    text_frame.text = 'Find the bullet slide layout'

    paragraph = text_frame.add_paragraph()
    paragraph.text = 'Use _TextFrame.text for first bullet'
    paragraph.level = 1

    paragraph = text_frame.add_paragraph()
    paragraph.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    paragraph.level = 2


def change_normal_slides(presentation):
    change_normal_slide(presentation.slides[1])
        


def create_pptx_file():
    #NOTE: You need to create the same number of slides in powerpoint for the template, I know it sucks but python-pptx cant copy slides.
    presentation = Presentation('aprg_template.pptx')

    change_title_slide(presentation)
    change_normal_slides(presentation)

    presentation.save('test.pptx')


if __name__ == "__main__":
    create_pptx_file()
    print('Done!')
